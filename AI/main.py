import os
import tempfile
import uuid
import base64
import zipfile
import tarfile
import io
import fitz  # PyMuPDF
import docx
import pptx
import boto3
import pinecone
import mimetypes
import requests
from dotenv import load_dotenv
from sentence_transformers import SentenceTransformer
from openai import OpenAI
import textwrap

load_dotenv()

# Load environment variables
pinecone_api_key = os.getenv("PINECONE_API_KEY")
pinecone_env = os.getenv("PINECONE_ENVIRONMENT")
pinecone_index_name = os.getenv("PINECONE_INDEX_NAME")
openai_api_key = os.getenv("OPENAI_API_KEY")
s3_bucket_name = os.getenv("S3_BUCKET_NAME")

# Initialize services
pinecone.init(api_key=pinecone_api_key, environment=pinecone_env)
index = pinecone.Index(pinecone_index_name)
client = OpenAI(api_key=openai_api_key)
s3 = boto3.client("s3")
embedding_model = SentenceTransformer("all-MiniLM-L6-v2")


def split_text(text, max_chunk_size=500):
    return textwrap.wrap(text, width=max_chunk_size, break_long_words=False, break_on_hyphens=False)


def extract_text_from_pdf(file_path):
    text = ""
    with fitz.open(file_path) as pdf:
        for page in pdf:
            text += page.get_text()
    return text


def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])


def extract_text_from_pptx(file_path):
    presentation = pptx.Presentation(file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def extract_text_from_zip(file_path):
    text = ""
    with zipfile.ZipFile(file_path, "r") as zip_ref:
        for zip_info in zip_ref.infolist():
            with zip_ref.open(zip_info.filename) as file:
                try:
                    content = file.read().decode("utf-8")
                    text += content + "\n"
                except:
                    continue
    return text


def extract_text_from_tar(file_path):
    text = ""
    with tarfile.open(file_path, "r:*") as tar:
        for member in tar.getmembers():
            if member.isfile():
                f = tar.extractfile(member)
                try:
                    content = f.read().decode("utf-8")
                    text += content + "\n"
                except:
                    continue
    return text


def extract_text_from_audio_api(file_path):
    """
    חלופה קלה יותר - שליחה ישירה ל-OpenAI API
    במקום להשתמש במוביפי ופאידאב
    """
    with open(file_path, "rb") as f:
        transcript = client.audio.transcriptions.create(
            model="whisper-1", 
            file=f
        )
    return transcript.text


def transcribe_image_with_ai(image_bytes):
    base64_image = base64.b64encode(image_bytes).decode("utf-8")
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "Extract all visible text from the image."},
            {
                "role": "user",
                "content": [
                    {"type": "image_url", 
                     "image_url": {"url": f"data:image/png;base64,{base64_image}"}}  
                ],
            },
        ],
    )
    return response.choices[0].message.content.strip()


def extract_text_from_image(file_path):
    with open(file_path, "rb") as image_file:
        image_bytes = image_file.read()
    return transcribe_image_with_ai(image_bytes)


def extract_text(file_path, file_type):
    """
    גרסה מאופטמת - הסרת תמיכה בוידאו כבד
    """
    if file_type == "pdf":
        return extract_text_from_pdf(file_path)
    elif file_type == "docx":
        return extract_text_from_docx(file_path)
    elif file_type == "pptx":
        return extract_text_from_pptx(file_path)
    elif file_type in ["zip", "tar"]:
        return extract_text_from_zip(file_path) if file_type == "zip" else extract_text_from_tar(file_path)
    elif file_type == "audio":
        return extract_text_from_audio_api(file_path)
    elif file_type == "image":
        return extract_text_from_image(file_path)
    else:
        # עבור וידאו - המר לאודיו בצד הלקוח או השתמש בשירות חיצוני
        return "Video processing not supported in lightweight version"


def get_file_type(file_path):
    mime_type, _ = mimetypes.guess_type(file_path)
    if mime_type:
        if "pdf" in mime_type:
            return "pdf"
        elif "word" in mime_type:
            return "docx"
        elif "presentation" in mime_type:
            return "pptx"
        elif "zip" in mime_type:
            return "zip"
        elif "tar" in mime_type:
            return "tar"
        elif "audio" in mime_type:
            return "audio"
        elif "video" in mime_type:
            return "video"
        elif "image" in mime_type:
            return "image"
    return "unknown"


def index_s3_file_for_user(user_id, s3_key):
    local_file_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
            s3.download_fileobj(s3_bucket_name, s3_key, temp_file)
            local_file_path = temp_file.name

        file_type = get_file_type(local_file_path)
        text = extract_text(local_file_path, file_type)
        chunks = split_text(text)

        for chunk in chunks:
            embedding = embedding_model.encode(chunk)
            index.upsert([
                {
                    "id": str(uuid.uuid4()),
                    "values": embedding.tolist(),
                    "metadata": {
                        "user_id": user_id,
                        "text": chunk,
                        "s3_key": s3_key
                    }
                }
            ])

    finally:
        if local_file_path and os.path.exists(local_file_path):
            os.remove(local_file_path)


def query_user_files(user_id, query, top_k=5):
    query_embedding = embedding_model.encode(query)
    results = index.query(
        vector=query_embedding.tolist(),
        top_k=top_k,
        include_metadata=True
    )

    filtered_results = []
    for match in results.matches:
        metadata = match.metadata
        if metadata.get("user_id") == user_id:
            filtered_results.append({
                "text": metadata.get("text"),
                "s3_key": metadata.get("s3_key"),
                "score": match.score,
                "s3_url": s3.generate_presigned_url(
                    "get_object",
                    Params={"Bucket": s3_bucket_name, "Key": metadata.get("s3_key")},
                    ExpiresIn=3600,
                )
            })
    return filtered_results